from __future__ import annotations

from pathlib import Path
from typing import Callable, Dict, List
from semantic_text_splitter import TextSplitter

max_characters = 1000
splitter = TextSplitter(max_characters)

# ──────────────────────────  Common helpers  ────────────────────────────

def _clean(text: str) -> str:
  return " ".join(text.split())

# ──────────────────────────  PDF  ───────────────────────────────────────

try:
  import pdfplumber  # type: ignore
except ImportError:  # pragma: no cover
  pdfplumber = None

def _parse_pdf(path: Path) -> str:
  out: List[str] = []
  with pdfplumber.open(path) as pdf:
    for i, page in enumerate(pdf.pages, 1):
      txt = page.extract_text() or ""
      out.append(f"[page {i}]\n{_clean(txt)}")
    return "\n\n".join(out)

# ──────────────────────────  PowerPoint  ────────────────────────────────

try:
  from pptx import Presentation  # type: ignore
except ImportError:  # pragma: no cover
  Presentation = None  # type: ignore

def _parse_pptx(path: Path) -> str:
  if Presentation is None:  # pragma: no cover
    raise RuntimeError("python-pptx is required for .pptx parsing; `pip install python-pptx`. ")
  prs = Presentation(path)
  slides: List[str] = []
  for idx, slide in enumerate(prs.slides, 1):
    texts = [
      shape.text.strip()
      for shape in slide.shapes
      if hasattr(shape, "text") and shape.text.strip()
    ]
    slides.append(f"[slide {idx}]\n{_clean(' '.join(texts))}")
  return "\n\n".join(slides)

# ──────────────────────────  Excel / CSV  ───────────────────────────────

try:
  import pandas as pd  # type: ignore
except ImportError:  # pragma: no cover
  pd = None  # type: ignore

def _parse_excel(path: Path) -> str:
  if pd is None:  # pragma: no cover
    raise RuntimeError("pandas + openpyxl are required for Excel parsing; `pip install pandas openpyxl`. ")

  if path.suffix.lower() == ".csv":
    df = pd.read_csv(path)
    return f"[csv]\n{df.to_markdown(index=False)}"

  all_sheets = pd.read_excel(path, sheet_name=None)
  out: List[str] = []
  for name, df in all_sheets.items():
    out.append(f"[sheet {name}]\n{df.to_markdown(index=False)}")
  return "\n\n".join(out)

# ──────────────────────────  Text / Markdown  ───────────────────────────

try:
  import markdown_it  # type: ignore
except ImportError:  # pragma: no cover
  markdown_it = None  # type: ignore

_MD = markdown_it.MarkdownIt("commonmark") if markdown_it else None  # type: ignore

def _parse_text(path: Path) -> str:
  raw = path.read_text(encoding="utf-8", errors="ignore")
  if path.suffix.lower() in {".md", ".markdown"} and _MD is not None:
    return _clean(_MD.render(raw))
  return _clean(raw)

# ──────────────────────────  Images (OCR)  ──────────────────────────────

try:
  from PIL import Image  # type: ignore
  import pytesseract  # type: ignore
except ImportError:  # pragma: no cover
  Image = None  # type: ignore
  pytesseract = None  # type: ignore

def _parse_image(path: Path, lang: str = "eng") -> str:
  if Image is None or pytesseract is None:  # pragma: no cover
    raise RuntimeError("pillow + pytesseract are required for OCR; `pip install pillow pytesseract`. ")
  img = Image.open(path)
  return _clean(pytesseract.image_to_string(img, lang=lang))

# ──────────────────────────  Dispatch table  ────────────────────────────

_EXTRACTOR_MAP: Dict[str, Callable[[Path], str]] = {
    ".pdf":  _parse_pdf,
    ".pptx": _parse_pptx,
    ".xlsx": _parse_excel,
    ".xls":  _parse_excel,
    ".csv":  _parse_excel,
    ".txt":  _parse_text,
    ".md":   _parse_text,
    ".markdown": _parse_text,
    ".png":  _parse_image,
    ".jpg":  _parse_image,
    ".jpeg": _parse_image,
    ".tiff": _parse_image,
}

# ──────────────────────────  Public API  ────────────────────────────────

def get_extractor(fileType: str | None = None) -> Callable[[Path], str]:
  ext = fileType.lower()
  func = _EXTRACTOR_MAP.get(ext)
  if func is None:
    raise ValueError(f"Unsupported file type: {ext}")
  return func
from pathlib import Path

def extract_chunks(path: str | Path) -> list[str]:
    path = Path(path)
    ext  = path.suffix.lower()            # e.g. ".pdf", ".pptx", ".csv", etc.
    extractor = _EXTRACTOR_MAP.get(ext)
    if extractor is None:
        raise ValueError(f"Unsupported file extension: {ext}")
    text = extractor(path)
    return splitter.chunks(text)